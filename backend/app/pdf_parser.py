"""
pdf_parser.py — Universal document text extractor + exam question parser.

Supports: PDF, DOCX, PPTX, XLSX/XLS/ODS, CSV/TSV, EPUB, HTML, images (OCR), plain text.

Install dependencies:
    pip install PyMuPDF pypdf python-docx python-pptx openpyxl pandas \
                pillow pytesseract ebooklib beautifulsoup4 xlrd odfpy nltk

Tesseract binary (for OCR):
    Ubuntu/Debian : sudo apt install tesseract-ocr
    Windows       : https://github.com/UB-Mannheim/tesseract/wiki
"""

import io
import logging
import os
import re
import base64
import string
from pathlib import Path
from typing import List, Dict, Any, Optional

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# Optional-dependency availability flags
# ─────────────────────────────────────────────────────────────────────────────

# PDF: PyMuPDF (primary)
try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    logger.warning("PyMuPDF (fitz) not installed — PDF extraction falls back to pypdf/PyPDF2")
    fitz = None
    FITZ_AVAILABLE = False

# PDF: pypdf (secondary)
try:
    from pypdf import PdfReader as _pypdf_PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    _pypdf_PdfReader = None
    PYPDF_AVAILABLE = False

# PDF: PyPDF2 (tertiary legacy fallback)
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PyPDF2 = None
    PYPDF2_AVAILABLE = False

# Image / OCR
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    logger.warning("Pillow not installed — image OCR unavailable")
    Image = None
    PIL_AVAILABLE = False

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    pytesseract = None
    TESSERACT_AVAILABLE = False

# DOCX
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    logger.warning("python-docx not installed (pip install python-docx)")
    DocxDocument = None
    DOCX_AVAILABLE = False

# PPTX
try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not installed (pip install python-pptx)")
    PptxPresentation = None
    PPTX_AVAILABLE = False

# XLSX
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    openpyxl = None
    OPENPYXL_AVAILABLE = False

# pandas (XLS / ODS / CSV fallback)
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    pd = None
    PANDAS_AVAILABLE = False

# EPUB
try:
    import ebooklib
    from ebooklib import epub as _epub
    from bs4 import BeautifulSoup as _BS
    EPUB_AVAILABLE = True
except ImportError:
    ebooklib = None
    _epub = None
    _BS = None
    EPUB_AVAILABLE = False

# NLTK keyword extraction
try:
    import nltk
    NLTK_AVAILABLE = True
except ImportError:
    nltk = None
    NLTK_AVAILABLE = False

# ─────────────────────────────────────────────────────────────────────────────
# NLTK lazy initialiser — downloads once per process, never per call
# ─────────────────────────────────────────────────────────────────────────────
_NLTK_READY: bool = False

def _ensure_nltk() -> bool:
    """
    Ensure NLTK punkt_tab + stopwords data are present.
    Downloads quietly on first call; subsequent calls are instant (flag check).
    Returns True if NLTK is usable, False otherwise.

    Windows note: nltk.data.find() raises OSError (not just LookupError) when
    a path is structurally valid but the file is missing — both are caught here.

    punkt vs punkt_tab: NLTK 3.8+ uses punkt_tab; older versions use punkt.
    We try punkt_tab first and fall back to punkt so both versions work.
    """
    global _NLTK_READY
    if _NLTK_READY:
        return True
    if not NLTK_AVAILABLE:
        return False

    # Try punkt_tab (NLTK 3.8+) first, fall back to punkt (older NLTK)
    tokenizer_ready = False
    for resource, data_path in [
        ("punkt_tab", "tokenizers/punkt_tab"),
        ("punkt",     "tokenizers/punkt"),
    ]:
        try:
            nltk.data.find(data_path)
            tokenizer_ready = True
            break  # found one — no need to try the other
        except (LookupError, OSError):
            try:
                nltk.download(resource, quiet=True)
                tokenizer_ready = True
                break
            except Exception as e:
                logger.debug("NLTK download skipped for '%s': %s", resource, e)
                continue  # try the next tokenizer variant

    if not tokenizer_ready:
        logger.warning("NLTK tokenizer unavailable — keyword extraction will use fallback")
        return False

    # Stopwords
    try:
        nltk.data.find("corpora/stopwords")
    except (LookupError, OSError):
        try:
            nltk.download("stopwords", quiet=True)
        except Exception as e:
            logger.warning("NLTK stopwords download failed: %s", e)
            return False

    _NLTK_READY = True
    return True


# ─────────────────────────────────────────────────────────────────────────────
# Main class
# ─────────────────────────────────────────────────────────────────────────────

class PDFParser:
    """
    Universal document parser — PDF, DOCX, PPTX, XLSX, CSV, EPUB, images, plain text.

    All public methods are @staticmethod so the class can be used without
    instantiation.  Entry point for most callers: PDFParser.extract_text(path).
    """

    # =========================================================================
    # PDF
    # =========================================================================

    @staticmethod
    def extract_text_from_pdf(pdf_path: str, ocr: bool = False) -> str:
        """
        Extract text from a PDF file.

        Fallback chain: PyMuPDF → pypdf → PyPDF2.
        If ocr=True, pages without a text layer are rasterised and passed
        through Tesseract (requires Pillow + pytesseract + Tesseract binary).
        Each page is labelled [Page N] for downstream context.
        """
        # 1. PyMuPDF
        if FITZ_AVAILABLE:
            try:
                with fitz.open(pdf_path) as doc:
                    return PDFParser._pages_from_fitz(doc, ocr)
            except Exception as exc:
                logger.warning("PyMuPDF failed on %s, trying pypdf: %s", pdf_path, exc)

        # 2. pypdf
        if PYPDF_AVAILABLE:
            try:
                reader = _pypdf_PdfReader(pdf_path)
                return PDFParser._pages_from_pypdf(reader)
            except Exception as exc:
                logger.warning("pypdf failed on %s, trying PyPDF2: %s", pdf_path, exc)

        # 3. PyPDF2
        if PYPDF2_AVAILABLE:
            try:
                with open(pdf_path, "rb") as fh:
                    reader = PyPDF2.PdfReader(fh)
                    return PDFParser._pages_from_pypdf2(reader)
            except Exception as exc:
                logger.error("PyPDF2 also failed on %s: %s", pdf_path, exc)
                raise

        raise RuntimeError(
            "No PDF library available. "
            "Install one: pip install PyMuPDF  OR  pip install pypdf"
        )

    @staticmethod
    def extract_text_from_pdf_bytes(pdf_bytes: bytes, ocr: bool = False) -> str:
        """
        Extract text from raw PDF bytes (no temp file written).
        PyMuPDF and pypdf both support in-memory streams directly.
        """
        # 1. PyMuPDF
        if FITZ_AVAILABLE:
            try:
                with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                    return PDFParser._pages_from_fitz(doc, ocr)
            except Exception as exc:
                logger.warning("PyMuPDF bytes extraction failed, trying pypdf: %s", exc)

        # 2. pypdf
        if PYPDF_AVAILABLE:
            try:
                reader = _pypdf_PdfReader(io.BytesIO(pdf_bytes))
                return PDFParser._pages_from_pypdf(reader)
            except Exception as exc:
                logger.warning("pypdf bytes extraction failed, trying PyPDF2: %s", exc)

        # 3. PyPDF2
        if PYPDF2_AVAILABLE:
            try:
                reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
                return PDFParser._pages_from_pypdf2(reader)
            except Exception as exc:
                logger.error("All in-memory PDF extraction methods failed: %s", exc)
                raise

        raise RuntimeError(
            "No PDF library available. "
            "Install one: pip install PyMuPDF  OR  pip install pypdf"
        )

    @staticmethod
    def extract_metadata_from_pdf(pdf_path: str) -> Dict[str, Any]:
        """Return a dict of PDF metadata (title, author, page count, …)."""
        if not FITZ_AVAILABLE:
            logger.warning("PyMuPDF not available — returning empty metadata")
            return {}
        try:
            with fitz.open(pdf_path) as doc:
                meta = doc.metadata
                return {
                    "title":             meta.get("title", ""),
                    "author":            meta.get("author", ""),
                    "subject":           meta.get("subject", ""),
                    "creator":           meta.get("creator", ""),
                    "producer":          meta.get("producer", ""),
                    "creation_date":     meta.get("creationDate", ""),
                    "modification_date": meta.get("modDate", ""),
                    "pages":             len(doc),
                    "encrypted":         doc.is_encrypted,
                }
        except Exception as exc:
            logger.error("Error extracting PDF metadata: %s", exc)
            return {}

    @staticmethod
    def extract_images_from_pdf(pdf_path: str) -> List[Dict[str, Any]]:
        """
        Extract all raster images embedded in a PDF as base64-encoded PNGs.
        CMYK images are converted to RGB before encoding (not silently skipped).
        """
        if not FITZ_AVAILABLE:
            logger.warning("PyMuPDF not available — skipping image extraction")
            return []

        images: List[Dict[str, Any]] = []
        try:
            with fitz.open(pdf_path) as doc:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    for img_index, img in enumerate(page.get_images(full=True)):
                        xref = img[0]
                        pix = None
                        try:
                            pix = fitz.Pixmap(doc, xref)
                            if pix.n >= 5:  # CMYK or CMYK+alpha
                                logger.debug(
                                    "Converting CMYK image xref=%d page=%d to RGB",
                                    xref, page_num,
                                )
                                pix = fitz.Pixmap(fitz.csRGB, pix)
                            img_data = pix.tobytes(output="png")
                            images.append({
                                "page":       page_num,
                                "index":      img_index,
                                "width":      pix.width,
                                "height":     pix.height,
                                "image_data": base64.b64encode(img_data).decode("utf-8"),
                                "colorspace": "RGB",
                            })
                        except Exception as img_exc:
                            logger.warning(
                                "Skipping image xref=%d page=%d: %s",
                                xref, page_num, img_exc,
                            )
                        finally:
                            pix = None  # explicit Pixmap release
        except Exception as exc:
            logger.error("Error extracting images from PDF: %s", exc)
        return images

    # ── Private PDF helpers ───────────────────────────────────────────────────

    @staticmethod
    def _pages_from_fitz(doc, ocr: bool) -> str:
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                pages.append(f"[Page {i+1}]\n{text.strip()}")
            elif ocr:
                pages.append(f"[Page {i+1} — OCR]\n{PDFParser._ocr_fitz_page(doc, i)}")
            else:
                pages.append(f"[Page {i+1}] (no text layer; pass ocr=True to extract)")
        return "\n\n".join(pages)

    @staticmethod
    def _pages_from_pypdf(reader) -> str:
        pages = []
        for i, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip()
            pages.append(f"[Page {i+1}]\n{text}" if text else f"[Page {i+1}] (no text layer)")
        return "\n\n".join(pages)

    @staticmethod
    def _pages_from_pypdf2(reader) -> str:
        pages = []
        for i, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip()
            pages.append(f"[Page {i+1}]\n{text}" if text else f"[Page {i+1}] (no text layer)")
        return "\n\n".join(pages)

    @staticmethod
    def _ocr_fitz_page(doc, page_index: int) -> str:
        """Rasterise a fitz page and run Tesseract OCR."""
        if not TESSERACT_AVAILABLE or not PIL_AVAILABLE:
            return "(OCR unavailable — install pytesseract and Pillow)"
        try:
            page = doc.load_page(page_index)
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            return pytesseract.image_to_string(img).strip() or "(no text detected by OCR)"
        except Exception as exc:
            return f"(OCR failed: {exc})"

    # =========================================================================
    # DOCX
    # =========================================================================

    @staticmethod
    def extract_text_from_docx(docx_path: str) -> str:
        """Extract text + table content from a DOCX file."""
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx is not installed. Run: pip install python-docx")
        try:
            doc = DocxDocument(docx_path)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells)
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)
        except Exception as exc:
            logger.error("DOCX extraction failed: %s", exc)
            raise

    # =========================================================================
    # PPTX
    # =========================================================================

    @staticmethod
    def extract_text_from_pptx(pptx_path: str) -> str:
        """Extract per-slide text from a PPTX file."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")
        try:
            prs = PptxPresentation(pptx_path)
            slides: List[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = [
                    para.text.strip()
                    for shape in slide.shapes
                    if shape.has_text_frame
                    for para in shape.text_frame.paragraphs
                    if para.text.strip()
                ]
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except Exception as exc:
            logger.error("PPTX extraction failed: %s", exc)
            raise

    # =========================================================================
    # XLSX / XLS / ODS
    # =========================================================================

    @staticmethod
    def extract_text_from_xlsx(xlsx_path: str) -> str:
        """
        Extract cell values from all sheets in an XLSX/XLSM workbook.
        Uses read_only mode to avoid loading the entire file into memory.
        Workbook is explicitly closed in a finally block to release the handle.
        """
        if not OPENPYXL_AVAILABLE:
            raise RuntimeError("openpyxl is not installed. Run: pip install openpyxl")
        wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
        try:
            sheets: List[str] = []
            for name in wb.sheetnames:
                ws = wb[name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join("" if c is None else str(c) for c in row)
                    if row_str.strip(" |"):
                        rows.append(row_str)
                if rows:
                    sheets.append(f"[Sheet: {name}]\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        finally:
            wb.close()  # always release the file handle

    @staticmethod
    def extract_text_from_xls(xls_path: str) -> str:
        """Extract text from a legacy .xls file using pandas + xlrd."""
        if not PANDAS_AVAILABLE:
            raise RuntimeError("pandas is not installed. Run: pip install pandas xlrd")
        xl = pd.ExcelFile(xls_path, engine="xlrd")
        parts = [
            f"[Sheet: {sheet}]\n{xl.parse(sheet).to_string(index=False)}"
            for sheet in xl.sheet_names
        ]
        return "\n\n".join(parts)

    @staticmethod
    def extract_text_from_ods(ods_path: str) -> str:
        """Extract text from an ODS file using pandas + odfpy."""
        if not PANDAS_AVAILABLE:
            raise RuntimeError("pandas is not installed. Run: pip install pandas odfpy")
        xl = pd.ExcelFile(ods_path, engine="odf")
        parts = [
            f"[Sheet: {sheet}]\n{xl.parse(sheet).to_string(index=False)}"
            for sheet in xl.sheet_names
        ]
        return "\n\n".join(parts)

    @staticmethod
    def extract_text_from_csv(csv_path: str) -> str:
        """Extract text from a CSV or TSV file."""
        if not PANDAS_AVAILABLE:
            return PDFParser.extract_text_from_txt(csv_path)
        sep = "\t" if csv_path.lower().endswith(".tsv") else ","
        df = pd.read_csv(csv_path, sep=sep)
        return df.to_string(index=False)

    # =========================================================================
    # EPUB
    # =========================================================================

    @staticmethod
    def extract_text_from_epub(epub_path: str) -> str:
        """Extract chapter text from an EPUB file."""
        if not EPUB_AVAILABLE:
            raise RuntimeError(
                "ebooklib/beautifulsoup4 not installed. "
                "Run: pip install ebooklib beautifulsoup4"
            )
        book = _epub.read_epub(epub_path)
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = _BS(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n").strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts)

    # =========================================================================
    # Plain text / HTML / images
    # =========================================================================

    @staticmethod
    def extract_text_from_txt(txt_path: str) -> str:
        """Read a plain-text file, trying utf-8 → latin-1 → cp1252."""
        for enc in ("utf-8", "latin-1", "cp1252"):
            try:
                return Path(txt_path).read_text(encoding=enc)
            except UnicodeDecodeError:
                continue
        raise RuntimeError(f"Could not decode text file with any known encoding: {txt_path}")

    @staticmethod
    def extract_text_from_html(html_path: str) -> str:
        """Strip HTML tags and return plain text (falls back to raw read)."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(Path(html_path).read_bytes(), "html.parser")
            return soup.get_text(separator="\n").strip()
        except ImportError:
            return PDFParser.extract_text_from_txt(html_path)

    @staticmethod
    def extract_text_from_image(image_path: str) -> str:
        """Run Tesseract OCR on a raster image file."""
        if not TESSERACT_AVAILABLE:
            raise RuntimeError(
                "pytesseract is not installed. "
                "Run: pip install pytesseract  (and install the Tesseract binary)"
            )
        if not PIL_AVAILABLE:
            raise RuntimeError("Pillow is not installed. Run: pip install Pillow")
        img = Image.open(image_path)
        return pytesseract.image_to_string(img).strip() or "(no text detected by OCR)"

    # =========================================================================
    # Universal dispatcher
    # =========================================================================

    @staticmethod
    def extract_text(file_path: str, ocr: bool = False) -> str:
        """
        Dispatch to the correct extractor based on file extension.

        Strips query-string fragments (e.g. 'file.docx?token=x') before
        determining the extension so the lookup never fails on a '?' suffix.

        Supported extensions:
            pdf, docx, doc, pptx, ppt, xlsx, xlsm, xls, ods, csv, tsv,
            epub, txt, md, log, rst, json, xml, html, htm,
            jpg, jpeg, png, bmp, tiff, tif, webp
        """
        clean_path = file_path.split("?")[0].rstrip()
        _, raw_ext = os.path.splitext(clean_path)
        ext = raw_ext.lower().lstrip(".")

        _IMAGE_EXTS = {"jpg", "jpeg", "png", "bmp", "tiff", "tif", "webp"}
        _TEXT_EXTS  = {"txt", "md", "log", "rst", "json", "xml"}

        if ext == "pdf":
            return PDFParser.extract_text_from_pdf(file_path, ocr=ocr)
        if ext in ("docx", "doc"):
            return PDFParser.extract_text_from_docx(file_path)
        if ext in ("pptx", "ppt"):
            return PDFParser.extract_text_from_pptx(file_path)
        if ext in ("xlsx", "xlsm"):
            return PDFParser.extract_text_from_xlsx(file_path)
        if ext == "xls":
            return PDFParser.extract_text_from_xls(file_path)
        if ext == "ods":
            return PDFParser.extract_text_from_ods(file_path)
        if ext in ("csv", "tsv"):
            return PDFParser.extract_text_from_csv(file_path)
        if ext == "epub":
            return PDFParser.extract_text_from_epub(file_path)
        if ext in ("html", "htm"):
            return PDFParser.extract_text_from_html(file_path)
        if ext in _TEXT_EXTS:
            return PDFParser.extract_text_from_txt(file_path)
        if ext in _IMAGE_EXTS:
            return PDFParser.extract_text_from_image(file_path)

        supported = sorted(
            {"pdf", "docx", "doc", "pptx", "ppt", "xlsx", "xlsm", "xls", "ods",
             "csv", "tsv", "epub", "html", "htm"} | _TEXT_EXTS | _IMAGE_EXTS
        )
        raise ValueError(
            f"Unsupported file type: .{ext}. "
            f"Supported: {', '.join(supported)}"
        )

    # =========================================================================
    # Exam question parser  (PrepIQ domain logic — keep separate from I/O above)
    # =========================================================================

    @staticmethod
    def parse_questions_from_text(text: str) -> List[Dict[str, Any]]:
        """
        Parse exam questions from extracted document text.

        Handles the wide variety of numbering formats found in Indian university
        papers (BPUT, VTU, Anna Univ, etc.):
          - Numbered   : "1.", "1)", "(1)", "Q1.", "Q.1", "Q.No.1", "Question 1"
          - Lettered   : "a)", "(a)", "A."
          - Roman      : "i)", "(i)", "ii.", "III)"
          - Marks lines: "(5 marks)", "[10M]", "5M" anywhere on the line
          - Verb cues  : lines starting with "Explain", "Define", "Prove", …
          - Fallback   : every substantial line (≥20 chars) when nothing matches
        """
        MAX_TEXT = 300_000
        text = text[:MAX_TEXT]

        questions: List[Dict[str, Any]] = []
        seen_hashes: set = set()

        # Header / boilerplate prefixes to skip
        _SKIP_PREFIXES = (
            "page ", "reg no", "register", "roll no", "date:", "time:",
            "max marks", "maximum marks", "total marks", "duration",
            "instructions", "note:", "answer all", "answer any",
            "part a", "part b", "part c", "section a", "section b",
            "unit i", "unit ii", "unit iii", "unit iv", "unit v",
            "module ", "all questions carry",
        )

        def _add(q_text: str, q_number: int) -> None:
            q_text = q_text.strip()
            if len(q_text) < 8:
                return
            lower = q_text.lower()
            if any(lower.startswith(p) for p in _SKIP_PREFIXES):
                return
            key = hash(lower[:120])
            if key in seen_hashes:
                return
            seen_hashes.add(key)
            questions.append({
                "text":          q_text,
                "number":        q_number,
                "marks":         PDFParser._estimate_marks(q_text),
                "unit":          PDFParser._estimate_unit(q_text),   # None = unknown
                "question_type": PDFParser._classify_question_type(q_text),
                "difficulty":    PDFParser._estimate_difficulty(q_text),
                "keywords":      PDFParser._extract_keywords(q_text),
                "length":        len(q_text),
            })

        # Pre-compiled patterns (compiled once per call, not per line)
        _num = re.compile(
            r'^(?:(?:Q(?:uestion|n|\.No?\.?)?\.?\s*)?\(?(\d{1,3})\)?[.):\s]\s*'
            r'|\((\d{1,3})\)\s+)',
            re.IGNORECASE,
        )
        _letter     = re.compile(r'^\(?([a-hA-H])\)?[.)]\s+(.+)', re.IGNORECASE)
        _roman      = re.compile(r'^\(?([ivxIVX]{1,4})\)?[.)]\s+(.+)')
        _marks_line = re.compile(
            r'(?:\(\s*\d+\s*(?:marks?)?\s*\)'
            r'|\[\s*\d+\s*(?:marks?|M)\s*\]'
            r'|\b\d+\s*marks?\b'
            r'|\b\d+\s*M\b)',
            re.IGNORECASE,
        )
        _q_verbs = re.compile(
            r'\b(?:explain|describe|define|discuss|derive|prove|show|find|calculate'
            r'|evaluate|compare|contrast|differentiate|list|write|state|what|why'
            r'|how|when|where|which|illustrate|analyze|analyse|design|implement'
            r'|develop|construct|draw|sketch|outline|summarize|justify)\b',
            re.IGNORECASE,
        )

        lines = [ln.strip() for ln in text.splitlines()]
        # FIX: start at 1 so the first question number is never 0
        q_counter = 0
        pending:     str = ""
        pending_num: int = 0

        def flush_pending() -> None:
            nonlocal pending, pending_num
            if pending:
                _add(pending, pending_num)
            pending = ""
            pending_num = 0

        for line in lines:
            if not line:
                flush_pending()
                continue

            m = _num.match(line)
            if m:
                flush_pending()
                q_counter += 1
                # FIX: fall back to q_counter (already incremented), never 0
                raw_num = m.group(1) or m.group(2)
                num = int(raw_num) if raw_num is not None else q_counter
                rest = line[m.end():].strip()
                if rest:
                    pending     = rest
                    pending_num = num
                continue

            m = _letter.match(line)
            if m:
                flush_pending()
                q_counter += 1
                pending     = m.group(2).strip()
                pending_num = q_counter
                continue

            m = _roman.match(line)
            if m:
                flush_pending()
                q_counter += 1
                rest        = line[m.end():].strip()
                pending     = rest if rest else line
                pending_num = q_counter
                continue

            if _marks_line.search(line) and len(line) > 15:
                clean = _marks_line.sub("", line).strip().strip("()[]").strip()
                if len(clean) >= 8:
                    flush_pending()
                    q_counter += 1
                    pending     = clean
                    pending_num = q_counter
                    continue

            if len(line) >= 20 and (line.endswith("?") or _q_verbs.search(line)):
                if pending:
                    pending = pending + " " + line
                else:
                    flush_pending()
                    q_counter += 1
                    pending     = line
                    pending_num = q_counter
                continue

            # Continuation line: append to current pending question
            if pending and len(line) >= 5:
                pending = pending + " " + line

        flush_pending()

        # Fallback: treat every substantial line as its own question
        if not questions:
            q_counter = 0
            for line in lines:
                line = line.strip()
                if len(line) >= 20:
                    q_counter += 1
                    _add(line, q_counter)

        return questions

    # =========================================================================
    # Exam-question metadata helpers
    # =========================================================================

    @staticmethod
    def _estimate_marks(question_text: str) -> int:
        """
        Extract marks from question text via regex.
        Patterns handled: "(10 marks)", "[5M]", "10 marks", "(10)", "5 Marks".
        Returns 0 when no marks indicator is found (0 = unknown, not a default).
        """
        text = question_text.lower()
        m = re.search(r'\b(\d{1,3})\s*marks?\b', text)
        if m:
            return int(m.group(1))
        m = re.search(r'[\[(]\s*(\d{1,3})\s*m\s*[\])]', text)
        if m:
            return int(m.group(1))
        # Bare "(N)" at end of line — likely a marks indicator
        m = re.search(r'\(\s*(\d{1,3})\s*\)\s*$', text)
        if m:
            val = int(m.group(1))
            if 1 <= val <= 20:
                return val
        return 0

    @staticmethod
    def _estimate_unit(question_text: str) -> Optional[str]:
        """
        Extract unit/module reference from question text.
        Returns None when no reference is found — do NOT assign a default.
        """
        text = question_text.lower()
        m = re.search(r'\b(?:unit|module)\s*[-:]?\s*([ivxIVX\d]+)\b', text)
        if m:
            return f"Unit {m.group(1).upper()}"
        m = re.search(r'\bco\s*(\d)\b', text)
        if m:
            return f"CO{m.group(1)}"
        return None

    @staticmethod
    def _classify_question_type(question_text: str) -> str:
        """Classify question by the primary verb/keyword in the text."""
        lower = question_text.lower()
        if any(k in lower for k in ("prove", "show that", "demonstrate", "derive")):
            return "Proof/derivation"
        if any(k in lower for k in ("calculate", "find", "evaluate", "solve")):
            return "Calculation/problem"
        if any(k in lower for k in ("explain", "describe", "discuss", "what is")):
            return "Conceptual/explanation"
        if any(k in lower for k in ("define", "state")):
            return "Definition"
        if any(k in lower for k in ("compare", "contrast", "difference")):
            return "Comparison"
        return "Mixed/other"

    @staticmethod
    def _estimate_difficulty(question_text: str) -> str:
        """
        Estimate difficulty.
        Primary signal: marks value (most reliable).
          ≤3  → Easy | 4–7 → Medium | ≥8 → Hard
        Secondary: keyword scan when no marks are found.
        """
        marks = PDFParser._estimate_marks(question_text)
        if marks > 0:
            if marks <= 3:
                return "Easy"
            if marks <= 7:
                return "Medium"
            return "Hard"

        lower = question_text.lower()
        hard_words = {
            "prove", "derive", "theorem", "proof", "advanced", "complex",
            "challenging", "critical", "analyse", "analyze", "design",
            "implement", "develop", "construct",
        }
        easy_words = {
            "define", "list", "state", "name", "basic", "simple",
            "introductory", "what is", "identify",
        }
        hard_score = sum(1 for w in hard_words if w in lower)
        easy_score = sum(1 for w in easy_words if w in lower)
        if hard_score > easy_score:
            return "Hard"
        if easy_score > hard_score:
            return "Easy"
        return "Medium"

    @staticmethod
    def _extract_keywords(question_text: str) -> List[str]:
        """
        Extract up to 5 meaningful keywords from question text.

        Uses NLTK punkt tokeniser + stopword list when available.
        NLTK data is downloaded at most once per process (_ensure_nltk flag).
        Falls back to a simple split-and-filter if NLTK is unavailable.
        """
        if _ensure_nltk():
            try:
                from nltk.corpus import stopwords
                from nltk.tokenize import word_tokenize

                stop_words = set(stopwords.words("english"))
                tokens = word_tokenize(question_text)
                keywords = []
                for tok in tokens:
                    clean = tok.translate(str.maketrans("", "", string.punctuation)).lower()
                    if clean and clean not in stop_words and len(clean) > 2:
                        keywords.append(clean)
                return list(dict.fromkeys(keywords))[:5]  # preserve order, dedupe
            except Exception as exc:
                logger.debug("NLTK keyword extraction failed, using fallback: %s", exc)

        # Fallback: split + strip punctuation
        words = question_text.split()
        keywords = [w.strip('.,!?()[]{}"\'') for w in words if len(w) > 3]
        return list(dict.fromkeys(keywords))[:5]